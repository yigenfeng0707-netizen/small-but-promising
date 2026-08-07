#!/usr/bin/env python3
"""
Generate 安居智评 Agent PPT presentation using python-pptx.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
DARK_GREEN = RGBColor(0x2E, 0x7D, 0x32)
DEEP_BLUE = RGBColor(0x15, 0x65, 0xC0)
WARN_RED = RGBColor(0xC6, 0x28, 0x28)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_shape(slide, left, top, width, height, fill_color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_title_bar(slide, title, subtitle=""):
    add_shape(slide, 0, 0, SW, Inches(1.2), DARK_GREEN)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7), title, font_size=32, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.4), subtitle, font_size=16, color=RGBColor(0xA5, 0xD6, 0xA7))

def add_content_slide(title, bullets, bullets2=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_shape(slide, 0, 0, SW, Inches(1.2), DARK_GREEN)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7), title, font_size=28, bold=True, color=WHITE)

    if bullets:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_GRAY
            p.font.name = "Microsoft YaHei"
            p.space_after = Pt(8)

    if bullets2:
        txBox = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.5), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets2):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_GRAY
            p.font.name = "Microsoft YaHei"
            p.space_after = Pt(8)

    return slide

# ==================== Slide 1: Cover ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_GREEN)
add_text(slide, Inches(0.5), Inches(2), Inches(12), Inches(1.2), "安居智评 Agent", font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(3.2), Inches(12), Inches(0.8), "让每个家庭都看得见化学品风险", font_size=24, color=RGBColor(0xA5, 0xD6, 0xA7), alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(4.5), Inches(12), Inches(0.6), "2026「小有可为」AI 向善创新挑战赛 · 绿色发展赛道", font_size=16, color=RGBColor(0x81, 0xC7, 0x84), alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.4), "拍照 / 语音，10 秒评测家庭化学品安全", font_size=14, color=RGBColor(0x66, 0xBB, 0x6A), alignment=PP_ALIGN.CENTER)

# ==================== Slide 2: Pain Points ====================
add_content_slide("痛点 —— 家里最危险的地方，可能是水槽下", [
    "我国家庭化学品急性中毒事件中，儿童误服占比 > 40%，1-4 岁最高发",
    "农村独居老人因看不懂标签、子女不在身边，风险显著更高",
    "市面缺乏「拍一下就知道安不安全」的轻量工具",
    "老人面对「次氯酸钠」「烷基苯磺酸钠」无从判断危险",
    "出事时不知道正确处置（如误服洁厕剂不应催吐）",
])

# ==================== Slide 3: Solution ====================
add_content_slide("方案 —— 拍一下，10 秒知道安不安全", [
    "一句话方案：拍一张照 / 说一句话，得到针对这瓶化学品 + 这个家庭的完整安全评测与应急方案",
    "",
    "三种入口：",
    "  · 拍照评测 —— 低识字友好",
    "  · 语音提问 —— 老人/视障友好",
    "  · 文本输入 —— 二次提问",
    "",
    "全链路 6 步：识别 → 成分解析 → 风险评测 → 家庭画像 → 场景建议 → 应急指导",
])

# ==================== Slide 4: Architecture ====================
add_content_slide("技术架构", [
    "前端：React SPA → FastAPI 单端口托管",
    "",
    "编排层：",
    "  · 串联 6 Agent",
    "  · Step5 场景建议 ‖ Step6 应急指导 并行（省 30% 延迟）",
    "  · 单 Agent 失败 → 降级返回部分结果",
    "",
    "底层能力：",
    "  · 阿里云百炼 Qwen3（文本）+ Qwen-VL（多模态）",
    "  · MSDS RAG 知识库（216 条 + FAISS 语义检索）",
])

# ==================== Slide 5: 6 Agents ====================
add_content_slide("6 Agent 协作 —— 从识别到应急，端到端覆盖", [
    "Step 1  识别 Agent（Qwen-VL）",
    "       → 从图片识别化学品名/品牌/成分表",
    "",
    "Step 2  成分解析 Agent（Qwen3 + RAG）",
    "       → 成分名归一化 + 查 MSDS 知识库",
    "",
    "Step 3  风险评测 Agent",
    "       → 毒性/易燃/腐蚀/过敏/环境 5 维评分",
], [
    "Step 4  家庭画像 Agent",
    "       → 结合家庭成员动态调整风险等级",
    "",
    "Step 5  场景建议 Agent",
    "       → 存储 + 防护 + 绿色替代品",
    "",
    "Step 6  应急指导 Agent",
    "       → 误服/误触/泄漏处置 + 一键呼救",
])

# ==================== Slide 6: Agent Table ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "6 Agent 各司其职")
# Table
rows, cols = 7, 3
table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(12), Inches(5)).table
table.columns[0].width = Inches(2.5)
table.columns[1].width = Inches(5)
table.columns[2].width = Inches(4.5)
headers = ["Agent", "职责", "关键产出"]
data = [
    ["识别 Agent", "图像识别化学品名/品牌/成分", "名称 + 成分列表"],
    ["成分解析 Agent", "成分名归一化 + 查 MSDS", "标准化成分 + MSDS 命中"],
    ["风险评测 Agent", "5 维评分（毒/燃/腐/敏/环境）", "风险等级 + 评分依据"],
    ["家庭画像 Agent", "结合家庭画像个性化调整", "上调等级 + 脆弱人群提示"],
    ["场景建议 Agent", "存储 + 防护 + 绿色替代品", "三类建议"],
    ["应急指导 Agent", "误服/误触/泄漏处置", "应急步骤 + 一键呼救"],
]
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_GREEN
for r, row in enumerate(data):
    for c, val in enumerate(row):
        cell = table.cell(r+1, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_GRAY
        if r % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY

# ==================== Slide 7: Knowledge Base ====================
add_content_slide("知识库与 RAG：让评测有据可查", [
    "216 条 MSDS 数据，覆盖 6 大品类",
    "",
    "清洁剂（43）  消毒剂（32）  农药（37）",
    "药品（40）    化妆品（31）  其他（33）",
    "",
    "检索方式：",
    "  1. 精确匹配（成分名/别名/产品名）",
    "  2. 子串匹配",
    "  3. FAISS 语义向量检索（Embedding + 余弦相似度）",
    "  4. difflib 相似度兜底",
    "",
    "风险评测 Agent 必须优先引用 MSDS 命中片段",
])

# ==================== Slide 8: Public Welfare ====================
add_content_slide("公益价值：4 个真实落地场景", [
    "家庭日常评测",
    "  · C 端用户拍照即知家里这瓶东西安不安全",
    "",
    "社区安全普查",
    "  · 社区工作者批量上传，导出 PDF 报告",
    "  · 定位高风险家庭，精准上门提醒",
    "",
    "学校科普",
    "  · 评测结果转科普卡片，做「家里的化学课」",
    "",
    "独居老人关怀",
    "  · 志愿者上门帮老人评测，重点关注误服风险",
    "  · 对接壹基金/敬老院/街道办",
])

# ==================== Slide 9: Demo Screenshot ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Demo 实操 —— 拍照评测一瓶洁厕剂")
add_text(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(0.5), "此处粘贴 3 张 Demo 真实截图（首页 / 结果页上半 / 结果页下半+呼救按钮）", font_size=16, color=MID_GRAY)
add_shape(slide, Inches(1), Inches(2.2), Inches(11), Inches(4.5), LIGHT_GRAY)
add_text(slide, Inches(4), Inches(4), Inches(5), Inches(1), "Demo 截图占位", font_size=20, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# ==================== Slide 10: Innovation ====================
add_content_slide("5 个差异化创新 —— 为什么是安居智评", [
    "1. 绿色替代品推荐",
    "   唯一把「绿色低毒替代品」作为强约束输出的家庭化学品评测工具",
    "",
    "2. 家庭画像个性化",
    "   同一瓶化学品，含婴儿与仅成年人，风险等级完全不同",
    "",
    "3. 全链路 6 Agent",
    "   从识别到应急，端到端覆盖，非单点功能",
    "",
    "4. 多模态入口",
    "   拍照 + 语音 + 文本，照顾低识字/老人/视障",
    "",
    "5. 公益批量评测接口",
    "   预留 /api/batch-evaluate，直接服务社区/学校",
])

# ==================== Slide 11: Review Dimensions ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "评审四维度 —— 我们的设计一一对应")
rows, cols = 5, 3
table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(12), Inches(5)).table
table.columns[0].width = Inches(2.5)
table.columns[1].width = Inches(5)
table.columns[2].width = Inches(4.5)
headers = ["评审维度", "对应设计", "落地证据"]
data = [
    ["公益价值", "4 个落地场景 + 批量评测接口 + 多模态照顾弱势", "PARTNER_API.md、独居老人关怀"],
    ["技术可行性", "FastAPI + React + 百炼 Qwen + Docker，33 测试通过", "Dockerfile、自动部署、keepalive"],
    ["用户友好度", "拍照/语音/文本三入口 + 一键呼救 + 结果分块", "前端三入口、120/中毒热线按钮"],
    ["创新价值", "绿色替代品 + 家庭画像 + 全链路 + 多模态 + 批量", "5 个差异化创新点"],
]
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_GREEN
for r, row in enumerate(data):
    for c, val in enumerate(row):
        cell = table.cell(r+1, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_GRAY
        if r % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY

# ==================== Slide 12: Tech Stack ====================
add_content_slide("技术栈 —— 主流稳定，可复现可部署", [
    "前端：React 18 + Vite + Axios + react-router-dom",
    "后端：FastAPI + Uvicorn + Pydantic + asyncio",
    "大模型：阿里云百炼 Qwen3 + Qwen-VL + text-embedding-v3",
    "知识库：216 条 MSDS + FAISS 向量检索 + 精确/模糊匹配混合 RAG",
    "部署：Docker + 魔搭创空间（GitHub Actions 自动部署）",
    "测试：pytest，33 个测试通过",
    "运维：keepalive 心跳保活 + 健康检查 + 自动重启",
])

# ==================== Slide 13: Roadmap ====================
add_content_slide("后续规划 —— 从初赛到长期愿景", [
    "8 月      初赛提交：PDF + PPT + Demo 视频",
    "8-9 月    决赛打磨：RAG 升级，MSDS 扩至 200+ 条",
    "9-10 月   公益试点：联系 1-2 个社区/学校做真实试点",
    "2026 年底 接入智能家居设备做「被动风险监测」",
    "长期愿景  成为家庭安全领域的「家庭医生」",
    "开放 MSDS 数据众包，建立行业标杆",
])

# ==================== Slide 14: Team ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "团队 & 致谢")
add_text(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(0.6), "团队成员", font_size=20, bold=True, color=DARK_GREEN)
add_text(slide, Inches(0.5), Inches(2.5), Inches(12), Inches(2), "（填写成员姓名 + 角色 + 一句话简介）\n\n\n指导老师 / 合作机构（占位）", font_size=14, color=MID_GRAY)
add_text(slide, Inches(0.5), Inches(4.8), Inches(12), Inches(0.6), "致谢", font_size=20, bold=True, color=DARK_GREEN)
add_text(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(1.5), "阿里云百炼平台 1 亿 Token 补贴\n赛事组委会\n试点社区 / 学校 / 公益机构", font_size=14, color=MID_GRAY)

# ==================== Slide 15: Closing ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_GREEN)
add_text(slide, Inches(0.5), Inches(2.5), Inches(12), Inches(1.2), "让每个家庭都看得见化学品风险", font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(3.8), Inches(12), Inches(0.8), "安居智评 Agent · 谢谢评审", font_size=20, color=RGBColor(0xA5, 0xD6, 0xA7), alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(0.5), "在线 Demo：https://gsym236998-home-chem-safety-agent.ms.show", font_size=14, color=RGBColor(0x81, 0xC7, 0x84), alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(6.2), Inches(12), Inches(0.5), "代码仓库：https://github.com/yigenfeng0707-netizen/small-but-promising", font_size=14, color=RGBColor(0x81, 0xC7, 0x84), alignment=PP_ALIGN.CENTER)

# Save
output_path = r"D:\APPs\small_but_promising\安居智评_PPT_v1.pptx"
prs.save(output_path)
print(f"PPT saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
